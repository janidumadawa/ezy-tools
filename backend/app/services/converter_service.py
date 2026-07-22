import os
import time
from pathlib import Path
from PyPDF2 import PdfReader, PdfWriter
from PIL import Image
from ..config import DOWNLOADS_DIR

try:
    from docx2pdf import convert as docx2pdf_convert
except:
    docx2pdf_convert = None

try:
    from pdf2docx import Converter as Pdf2DocxConverter
except:
    Pdf2DocxConverter = None

class ConverterService:
    def __init__(self):
        self.download_dir = DOWNLOADS_DIR / "converted"
        self.download_dir.mkdir(parents=True, exist_ok=True)
    
    def _sanitize_filename(self, name: str) -> str:
        return "".join(c for c in name if c.isalnum() or c in (' ', '-', '_', '.')).rstrip()
    
    def word_to_pdf(self, file) -> dict:
        """Convert Word to PDF using LibreOffice (preserves all formatting)"""
        try:
            import subprocess
            
            original_name = Path(file.filename).stem
            safe_name = self._sanitize_filename(original_name)
            
            # Save uploaded file
            temp_docx = self.download_dir / f"{safe_name}.docx"
            with open(temp_docx, 'wb') as f:
                f.write(file.file.read())
            
            output_filename = f"{safe_name}.pdf"
            output_path = self.download_dir / output_filename
            
            # Try LibreOffice first (best quality, preserves formatting)
            try:
                subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', str(self.download_dir),
                    str(temp_docx)
                ], timeout=60, check=True, capture_output=True)
                
                # LibreOffice creates file with same name but .pdf extension
                expected_pdf = self.download_dir / f"{safe_name}.pdf"
                if expected_pdf.exists():
                    os.unlink(temp_docx)
                    return {
                        'success': True,
                        'data': {
                            'filename': output_filename,
                            'download_url': f'/api/converter/file/{output_filename}',
                            'message': f'{original_name}.pdf ready'
                        }
                    }
            except:
                pass
            
            # Fallback: Use python-docx + reportlab (text only, minimal formatting)
            try:
                from docx import Document
                from reportlab.lib.pagesizes import A4
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
                from reportlab.lib.styles import getSampleStyleSheet
                from reportlab.lib import colors
                from reportlab.lib.enums import TA_CENTER, TA_LEFT
                
                doc = Document(str(temp_docx))
                pdf_doc = SimpleDocTemplate(str(output_path), pagesize=A4)
                styles = getSampleStyleSheet()
                story = []
                
                for para in doc.paragraphs:
                    if para.text.strip():
                        text = para.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        
                        # Check for heading styles
                        style_name = para.style.name if para.style else ''
                        if 'Heading 1' in style_name:
                            style = styles['Heading1']
                        elif 'Heading 2' in style_name:
                            style = styles['Heading2']
                        elif 'Heading 3' in style_name:
                            style = styles['Heading3']
                        else:
                            style = styles['Normal']
                        
                        # Check alignment
                        if para.alignment == 1:  # Center
                            style.alignment = TA_CENTER
                        else:
                            style.alignment = TA_LEFT
                        
                        story.append(Paragraph(text, style))
                        story.append(Spacer(1, 6))
                
                # Handle tables
                for table in doc.tables:
                    data = []
                    for row in table.rows:
                        row_data = [cell.text.replace('&', '&amp;') for cell in row.cells]
                        data.append(row_data)
                    
                    if data:
                        t = Table(data)
                        t.setStyle(TableStyle([
                            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#884ab2')),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                            ('FONTSIZE', (0, 0), (-1, -1), 9),
                            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                            ('PADDING', (0, 0), (-1, -1), 6),
                        ]))
                        story.append(t)
                        story.append(Spacer(1, 12))
                
                pdf_doc.build(story)
                os.unlink(temp_docx)
                
                return {
                    'success': True,
                    'data': {
                        'filename': output_filename,
                        'download_url': f'/api/converter/file/{output_filename}',
                        'message': f'{original_name}.pdf ready'
                    }
                }
            except:
                os.unlink(temp_docx)
                raise Exception("Neither LibreOffice nor python-docx conversion worked")
                
        except Exception as e:
            return {'success': False, 'error': str(e)}

     
    def pdf_to_word(self, file) -> dict:
        try:
            if not Pdf2DocxConverter:
                return {'success': False, 'error': 'pdf2docx not installed'}
            original_name = Path(file.filename).stem
            safe_name = self._sanitize_filename(original_name)
            temp_pdf = self.download_dir / f"{safe_name}.pdf"
            with open(temp_pdf, 'wb') as f:
                f.write(file.file.read())
            output_filename = f"{safe_name}.docx"
            output_path = self.download_dir / output_filename
            cv = Pdf2DocxConverter(str(temp_pdf))
            cv.convert(str(output_path))
            cv.close()
            os.unlink(temp_pdf)
            return {'success': True, 'data': {'filename': output_filename, 'download_url': f'/api/converter/file/{output_filename}', 'message': f'{original_name}.docx ready'}}
        except Exception as e:
            return {'success': False, 'error': str(e)}
    
    def image_to_pdf(self, file) -> dict:
        try:
            original_name = Path(file.filename).stem
            safe_name = self._sanitize_filename(original_name)
            img = Image.open(file.file)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            output_filename = f"{safe_name}.pdf"
            output_path = self.download_dir / output_filename
            img.save(str(output_path), 'PDF')
            return {'success': True, 'data': {'filename': output_filename, 'download_url': f'/api/converter/file/{output_filename}', 'message': f'{original_name}.pdf ready'}}
        except Exception as e:
            return {'success': False, 'error': str(e)}

    def convert_image_format(self, file, target_format: str = 'png') -> dict:
        try:
            original_name = Path(file.filename).stem
            safe_name = self._sanitize_filename(original_name)
            img = Image.open(file.file)
            if target_format == 'jpg' and img.mode != 'RGB':
                img = img.convert('RGB')
            output_filename = f"{safe_name}.{target_format}"
            output_path = self.download_dir / output_filename
            img.save(str(output_path), target_format.upper(), quality=90)
            return {'success': True, 'data': {'filename': output_filename, 'download_url': f'/api/converter/file/{output_filename}', 'message': f'{original_name}.{target_format} ready'}}
        except Exception as e:
            return {'success': False, 'error': str(e)}

    def excel_to_pdf(self, file) -> dict:
        try:
            import openpyxl
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
            from reportlab.lib import colors
            original_name = Path(file.filename).stem
            safe_name = self._sanitize_filename(original_name)
            wb = openpyxl.load_workbook(file.file)
            ws = wb.active
            data = []
            for row in ws.iter_rows(values_only=True):
                data.append([str(cell) if cell is not None else '' for cell in row])
            output_filename = f"{safe_name}.pdf"
            output_path = self.download_dir / output_filename
            doc = SimpleDocTemplate(str(output_path), pagesize=landscape(A4))
            table = Table(data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTSIZE', (0, 0), (-1, -1), 8),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ]))
            doc.build([table])
            return {'success': True, 'data': {'filename': output_filename, 'download_url': f'/api/converter/file/{output_filename}', 'message': f'{original_name}.pdf ready'}}
        except Exception as e:
            return {'success': False, 'error': str(e)}

    def text_to_pdf(self, file) -> dict:
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph
            from reportlab.lib.styles import getSampleStyleSheet
            original_name = Path(file.filename).stem
            safe_name = self._sanitize_filename(original_name)
            content = file.file.read().decode('utf-8', errors='ignore')
            output_filename = f"{safe_name}.pdf"
            output_path = self.download_dir / output_filename
            doc = SimpleDocTemplate(str(output_path), pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            for line in content.split('\n'):
                if line.strip():
                    story.append(Paragraph(line.replace('&', '&amp;').replace('<', '&lt;'), styles['Normal']))
            doc.build(story)
            return {'success': True, 'data': {'filename': output_filename, 'download_url': f'/api/converter/file/{output_filename}', 'message': f'{original_name}.pdf ready'}}
        except Exception as e:
            return {'success': False, 'error': str(e)}

    def pptx_to_pdf(self, file) -> dict:
        try:
            from pptx import Presentation
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            original_name = Path(file.filename).stem
            safe_name = self._sanitize_filename(original_name)
            temp_pptx = self.download_dir / f"{safe_name}.pptx"
            with open(temp_pptx, 'wb') as f:
                f.write(file.file.read())
            output_filename = f"{safe_name}.pdf"
            output_path = self.download_dir / output_filename
            prs = Presentation(str(temp_pptx))
            doc = SimpleDocTemplate(str(output_path), pagesize=landscape(A4))
            styles = getSampleStyleSheet()
            story = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        story.append(Paragraph(shape.text.replace('&', '&amp;').replace('<', '&lt;'), styles['Normal']))
                        story.append(Spacer(1, 12))
                story.append(Spacer(1, 20))
            doc.build(story)
            os.unlink(temp_pptx)
            return {'success': True, 'data': {'filename': output_filename, 'download_url': f'/api/converter/file/{output_filename}', 'message': f'{original_name}.pdf ready'}}
        except Exception as e:
            return {'success': False, 'error': str(e)}


    def csv_to_pdf(self, file) -> dict:
        try:
            import csv
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
            from reportlab.lib import colors
            original_name = Path(file.filename).stem
            safe_name = self._sanitize_filename(original_name)
            content = file.file.read().decode('utf-8', errors='ignore')
            reader = csv.reader(content.splitlines())
            data = [row for row in reader]
            if not data:
                return {'success': False, 'error': 'Empty CSV file'}
            output_filename = f"{safe_name}.pdf"
            output_path = self.download_dir / output_filename
            doc = SimpleDocTemplate(str(output_path), pagesize=landscape(A4))
            table = Table(data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#884ab2')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTSIZE', (0, 0), (-1, -1), 8),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ]))
            doc.build([table])
            return {'success': True, 'data': {'filename': output_filename, 'download_url': f'/api/converter/file/{output_filename}', 'message': f'{original_name}.pdf ready'}}
        except Exception as e:
            return {'success': False, 'error': str(e)}

    def merge_images_to_pdf(self, files) -> dict:
        try:
            images = []
            for file in files:
                img = Image.open(file.file)
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                images.append(img)
            output_filename = f"merged_images.pdf"
            output_path = self.download_dir / output_filename
            if len(images) == 1:
                images[0].save(str(output_path), 'PDF')
            else:
                images[0].save(str(output_path), 'PDF', save_all=True, append_images=images[1:])
            return {'success': True, 'data': {'filename': output_filename, 'download_url': f'/api/converter/file/{output_filename}', 'message': f'Merged {len(images)} images into PDF'}}
        except Exception as e:
            return {'success': False, 'error': str(e)}